import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # Brand Colors
    # -------------------------------------------------------------
    c_dark_navy = RGBColor(11, 25, 44)      # #0B192C
    c_cobalt = RGBColor(30, 62, 98)         # #1E3E62
    c_teal = RGBColor(0, 140, 140)          # #008C8C (Vibrant Tech Accent)
    c_light_bg = RGBColor(248, 250, 252)    # #F8FAFC
    c_white = RGBColor(255, 255, 255)       # #FFFFFF
    c_text_dark = RGBColor(30, 41, 59)      # #1E293B
    c_text_muted = RGBColor(100, 116, 139)  # #64748B
    c_light_border = RGBColor(226, 232, 240) # #E2E8F0
    
    blank_layout = prs.slide_layouts[6] # Blank layout for full freedom
    
    # Helper to set solid background color
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw clean rectangular cards
    def add_card(slide, left, top, width, height, fill_color, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background() # No border
        return shape

    # Helper to add modern typography titles
    def add_slide_header(slide, title, category="PORTFOLIO"):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = "Segoe UI"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_teal
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(10), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p_title = title_tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark_navy

    # Helper to add standard formatted text to shapes/textboxes
    def add_bullet(tf, bold_prefix, text_content, font_size=14, space_after=12):
        p = tf.add_paragraph()
        p.space_after = Pt(space_after)
        
        run_bold = p.add_run()
        run_bold.text = bold_prefix + "  "
        run_bold.font.name = "Segoe UI"
        run_bold.font.bold = True
        run_bold.font.size = Pt(font_size)
        run_bold.font.color.rgb = c_text_dark
        
        run_text = p.add_run()
        run_text.text = text_content
        run_text.font.name = "Segoe UI"
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = c_text_muted

    # =============================================================
    # SLIDE 1: Professional Cover Page (Dark Theme)
    # =============================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1, c_dark_navy)
    
    # Subtle glowing graphic element on left
    g1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    g1.fill.solid()
    g1.fill.fore_color.rgb = c_teal
    g1.line.fill.background()
    
    # Subtitle Category
    cov_cat = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(0.4))
    cov_cat.text_frame.word_wrap = True
    p_cov_cat = cov_cat.text_frame.paragraphs[0]
    p_cov_cat.text = "ENTERPRISE SOFTWARE ARCHITECT & FULL STACK DEVELOPER"
    p_cov_cat.font.name = "Segoe UI"
    p_cov_cat.font.size = Pt(11)
    p_cov_cat.font.bold = True
    p_cov_cat.font.color.rgb = c_teal
    
    # Title
    cov_title = s1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(1.8))
    tf_cov_title = cov_title.text_frame
    tf_cov_title.word_wrap = True
    p_title = tf_cov_title.paragraphs[0]
    p_title.text = "Full Stack & ERP Software\nDevelopment Portfolio"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    
    # Subtitle / Core Competencies
    cov_sub = s1.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(11), Inches(0.6))
    tf_cov_sub = cov_sub.text_frame
    tf_cov_sub.word_wrap = True
    p_cov_sub = tf_cov_sub.paragraphs[0]
    p_cov_sub.text = "Scalable Web Applications  |  Enterprise ERP Systems  |  High-Performance SaaS Platforms"
    p_cov_sub.font.name = "Segoe UI"
    p_cov_sub.font.size = Pt(16)
    p_cov_sub.font.color.rgb = RGBColor(165, 180, 252) # Light blueish
    
    # Developer metadata card
    add_card(s1, Inches(1.2), Inches(5.2), Inches(10.9), Inches(1.3), c_cobalt)
    meta_box = s1.shapes.add_textbox(Inches(1.4), Inches(5.35), Inches(10.5), Inches(1.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta1 = tf_meta.paragraphs[0]
    p_meta1.text = "Prawin Kumar N"
    p_meta1.font.name = "Segoe UI"
    p_meta1.font.bold = True
    p_meta1.font.size = Pt(18)
    p_meta1.font.color.rgb = c_white
    
    p_meta2 = tf_meta.add_paragraph()
    p_meta2.text = "Erode, Tamil Nadu  •  prawin.vercel.app  •  github.com/prawinkumar2k  •  linkedin.com/in/prawinkumar-n"
    p_meta2.font.name = "Segoe UI"
    p_meta2.font.size = Pt(13)
    p_meta2.font.color.rgb = RGBColor(200, 210, 230)

    # =============================================================
    # SLIDE 2: About Me (Light Theme - Split Layout)
    # =============================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2, c_light_bg)
    add_slide_header(s2, "About Me", "Professional Profile")
    
    # Left Dark Card (Core Philosophy)
    add_card(s2, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0), c_dark_navy)
    philosophy_box = s2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(3.6), Inches(4.4))
    tf_phil = philosophy_box.text_frame
    tf_phil.word_wrap = True
    
    p_phil_title = tf_phil.paragraphs[0]
    p_phil_title.text = "CORE MISSION"
    p_phil_title.font.name = "Segoe UI"
    p_phil_title.font.size = Pt(12)
    p_phil_title.font.bold = True
    p_phil_title.font.color.rgb = c_teal
    p_phil_title.space_after = Pt(20)
    
    p_phil_body = tf_phil.add_paragraph()
    p_phil_body.text = "Bridging advanced backend architecture with ultra-slick frontend execution to build enterprise-grade software that drives organizational efficiency."
    p_phil_body.font.name = "Segoe UI"
    p_phil_body.font.size = Pt(18)
    p_phil_body.font.color.rgb = c_white
    p_phil_body.space_after = Pt(30)
    
    p_phil_loc = tf_phil.add_paragraph()
    p_phil_loc.text = "📍 Based in Erode, Tamil Nadu\n🚀 Specialized in High-Volume Portals"
    p_phil_loc.font.name = "Segoe UI"
    p_phil_loc.font.size = Pt(13)
    p_phil_loc.font.color.rgb = RGBColor(165, 180, 252)

    # Right Content Cards
    add_card(s2, Inches(5.3), Inches(1.6), Inches(7.2), Inches(5.0), c_white, c_light_border)
    about_box = s2.shapes.add_textbox(Inches(5.6), Inches(1.9), Inches(6.6), Inches(4.4))
    tf_about = about_box.text_frame
    tf_about.word_wrap = True
    
    p_abt_h = tf_about.paragraphs[0]
    p_abt_h.text = "PROFESSIONAL SNAPSHOT"
    p_abt_h.font.name = "Segoe UI"
    p_abt_h.font.bold = True
    p_abt_h.font.size = Pt(13)
    p_abt_h.font.color.rgb = c_cobalt
    p_abt_h.space_after = Pt(15)
    
    add_bullet(tf_about, "✦ Scalable System Architect:", "Experienced in creating heavy-duty web systems with cluster mode capability, optimizing response times to under 50ms.", 14, 15)
    add_bullet(tf_about, "✦ Government Tech Excellence:", "Developed state-level public admission platforms running smoothly for over 50,000+ applicants.", 14, 15)
    add_bullet(tf_about, "✦ Industrial IoT Integration:", "Successfully built smart attendance systems syncing Arduino hardware/RFID logs to enterprise databases in real-time.", 14, 15)
    add_bullet(tf_about, "✦ Premium UI/UX Implementation:", "Fluent in crafting sleek client portals utilizing React 19, AG Grid, Framer Motion, and Tailwind CSS.", 14, 10)

    # =============================================================
    # SLIDE 3: Technical Skills (Light Theme - 5 Cards Grid)
    # =============================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3, c_light_bg)
    add_slide_header(s3, "Technical Expertise", "Skills Matrix")
    
    skills_data = [
        {"title": "FRONTEND", "items": ["React.js (19/Vite)", "Tailwind CSS", "Bootstrap 5", "Framer Motion", "AG Grid Enterprise", "Responsive UI/UX"]},
        {"title": "BACKEND", "items": ["Node.js (LTS)", "Express.js 5", "PHP", "RESTful API Dev", "Cluster Mode Multi-threading", "JWT/RBAC Auth"]},
        {"title": "DATABASES", "items": ["MySQL 8.0", "MongoDB", "SQL Server (T-SQL)", "Connection Pooling", "Indexing & Optimization", "Normalized Database Design"]},
        {"title": "DEVOPS & CLOUD", "items": ["Docker Containerization", "Kubernetes (K3s)", "Nginx Reverse Proxy", "Dedicated Server Admin", "Linux System Administration", "CI/CD Workflows"]},
        {"title": "TOOLS & LIBS", "items": ["Git / GitHub", "Puppeteer (PDF Generation)", "Winston Logging", "RFID/Arduino IoT", "Postman API Suite", "Vercel / Heroku"]}
    ]
    
    card_width = Inches(2.2)
    card_height = Inches(4.8)
    card_y = Inches(1.8)
    gap = Inches(0.18)
    start_x = Inches(0.8)
    
    for i, s_grp in enumerate(skills_data):
        cx = start_x + i * (card_width + gap)
        # Background card
        add_card(s3, cx, card_y, card_width, card_height, c_white, c_light_border)
        
        # Skill header card
        hdr_shape = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, card_y, card_width, Inches(0.6))
        hdr_shape.fill.solid()
        hdr_shape.fill.fore_color.rgb = c_cobalt
        hdr_shape.line.fill.background()
        
        # Title Text
        tf_hdr = hdr_shape.text_frame
        tf_hdr.word_wrap = True
        p_hdr = tf_hdr.paragraphs[0]
        p_hdr.text = s_grp["title"]
        p_hdr.font.name = "Segoe UI"
        p_hdr.font.size = Pt(12)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = c_white
        p_hdr.alignment = PP_ALIGN.CENTER
        
        # Skill details textbox
        sb_box = s3.shapes.add_textbox(cx + Inches(0.1), card_y + Inches(0.8), card_width - Inches(0.2), card_height - Inches(0.9))
        tf_sb = sb_box.text_frame
        tf_sb.word_wrap = True
        
        for item in s_grp["items"]:
            p_item = tf_sb.add_paragraph() if tf_sb.text else tf_sb.paragraphs[0]
            p_item.space_after = Pt(12)
            p_item.text = "• " + item
            p_item.font.name = "Segoe UI"
            p_item.font.size = Pt(12.5)
            p_item.font.color.rgb = c_text_dark

    # =============================================================
    # SLIDE 4: Government Admissions Portal (DOTE Tamil Nadu)
    # =============================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4, c_light_bg)
    add_slide_header(s4, "Government Admissions Portal (DOTE Tamil Nadu)", "Case Study 01")
    
    # Left Specs Card
    add_card(s4, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0), c_white, c_light_border)
    tf_dote = s4.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.2), Inches(4.4)).text_frame
    tf_dote.word_wrap = True
    
    p_dote_h = tf_dote.paragraphs[0]
    p_dote_h.text = "HIGH-VOLUME PUBLIC REGISTRATION PORTAL"
    p_dote_h.font.name = "Segoe UI"
    p_dote_h.font.bold = True
    p_dote_h.font.size = Pt(13)
    p_dote_h.font.color.rgb = c_teal
    p_dote_h.space_after = Pt(15)
    
    add_bullet(tf_dote, "👥 Massive User Volume:", "Successfully architected and deployed to handle 50,000+ active state-wide student admissions smoothly.", 13.5, 14)
    add_bullet(tf_dote, "⚡ Technology Stack:", "Powered by React.js on the frontend with a highly resilient Node.js Express clustered backend API.", 13.5, 14)
    add_bullet(tf_dote, "💳 Secure Payment Integration:", "Unified secure multi-bank payment gateway integration for real-time application fee collection.", 13.5, 14)
    add_bullet(tf_dote, "🔐 Strict Security Guardrails:", "Implemented strict JWT session handling and fine-grained Role-Based Access Control (RBAC) protecting sensitive files.", 13.5, 14)
    add_bullet(tf_dote, "🖥️ Dedicated Server Infrastructure:", "Engineered custom server deployment with Nginx routing, SSL certificates, load balancing, and active server health monitors.", 13.5, 5)

    # Right Image Showcase
    img_path_dote = "c:\\Users\\Hp\\Documents\\HMS\\server_hms\\assets\\dote_dashboard.png"
    if os.path.exists(img_path_dote):
        # Card border/shadow for image
        add_card(s4, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), c_white, c_light_border)
        s4.shapes.add_picture(img_path_dote, Inches(7.0), Inches(1.75), Inches(5.4), Inches(4.7))
    else:
        # Placeholder shape if image not available
        add_card(s4, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), c_cobalt)

    # =============================================================
    # SLIDE 5: Campus ERP System
    # =============================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5, c_light_bg)
    add_slide_header(s5, "Enterprise Campus ERP System", "Case Study 02")
    
    # Left Specifications
    add_card(s5, Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.0), c_white, c_light_border)
    tf_erp = s5.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.6), Inches(4.4)).text_frame
    tf_erp.word_wrap = True
    
    p_erp_h = tf_erp.paragraphs[0]
    p_erp_h.text = "INTEGRATED CAMPUS RESOURCE & ACADEMIC ENGINE"
    p_erp_h.font.name = "Segoe UI"
    p_erp_h.font.bold = True
    p_erp_h.font.size = Pt(13)
    p_erp_h.font.color.rgb = c_cobalt
    p_erp_h.space_after = Pt(15)
    
    add_bullet(tf_erp, "👨‍🎓 Full Student Lifecycle:", "Unified student record keeping spanning initial admission, profile administration, grading, and transfer certificates.", 13.5, 14)
    add_bullet(tf_erp, "📅 Real-time Attendance & eMAR:", "Automated attendance charting with daily tracking and instant SMS/Email notifications to stakeholders.", 13.5, 14)
    add_bullet(tf_erp, "📊 Academic Reporting Engine:", "Powerful dashboard generating detailed analytical academic reports, test scores, class performance matrices, and custom GPA cards.", 13.5, 14)
    add_bullet(tf_erp, "🗄️ SQL Server Database Core:", "Heavy database structure utilizing optimized stored procedures, indexed views, and optimized relational schemas.", 13.5, 14)
    add_bullet(tf_erp, "🎨 Premium Bootstrap UI Console:", "Highly visual administration panels with clean corporate styling, mobile responsiveness, and dynamic tables.", 13.5, 5)

    # Right visual block
    add_card(s5, Inches(7.3), Inches(1.6), Inches(5.2), Inches(5.0), c_dark_navy)
    tf_erp_viz = s5.shapes.add_textbox(Inches(7.6), Inches(2.0), Inches(4.6), Inches(4.2)).text_frame
    tf_erp_viz.word_wrap = True
    p_viz_t = tf_erp_viz.paragraphs[0]
    p_viz_t.text = "ERP SYSTEM SCHEMATICS"
    p_viz_t.font.name = "Segoe UI"
    p_viz_t.font.size = Pt(12)
    p_viz_t.font.bold = True
    p_viz_t.font.color.rgb = c_teal
    p_viz_t.space_after = Pt(25)
    
    p_viz_b = tf_erp_viz.add_paragraph()
    p_viz_b.text = "Integrated Modules:\n\n✔️ Admin Console\n✔️ Student Portal\n✔️ Faculty Station\n✔️ Finance & Fees Desk\n✔️ Library Catalog\n\nDatabase: SQL Server Core\nFrontend: Bootstrap 5 Dashboard"
    p_viz_b.font.name = "Segoe UI"
    p_viz_b.font.size = Pt(16)
    p_viz_b.font.color.rgb = c_white

    # =============================================================
    # SLIDE 6: Hospital Management System
    # =============================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6, c_light_bg)
    add_slide_header(s6, "Lifeline Hospital Management System (HMS)", "Case Study 03")
    
    # Left Specs
    add_card(s6, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0), c_white, c_light_border)
    tf_hms = s6.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.2), Inches(4.4)).text_frame
    tf_hms.word_wrap = True
    
    p_hms_h = tf_hms.paragraphs[0]
    p_hms_h.text = "CLINICAL & BILLING PATIENT LIFECYCLE MANAGEMENT"
    p_hms_h.font.name = "Segoe UI"
    p_hms_h.font.bold = True
    p_hms_h.font.size = Pt(13)
    p_hms_h.font.color.rgb = c_teal
    p_hms_h.space_after = Pt(15)
    
    add_bullet(tf_hms, "🏥 Clinical PLCM Flow:", "Tracks complete patient journey from Registration (OPD), Admission (IPD), Ward assignment, Nursing charts, down to discharge.", 13.5, 14)
    add_bullet(tf_hms, "👨‍⚕️ Role-Based Medical Stations:", "Dedicated clinical workstations for Admin, Doctor, Receptionist, Nurse, Pharmacist, Lab Technician, and Billing clerk.", 13.5, 14)
    add_bullet(tf_hms, "🔬 Pharmacy & Lab Integration:", "Dynamic pharmacy stock inventory control and digital laboratory test reports generation via Puppeteer PDF renderer.", 13.5, 14)
    add_bullet(tf_hms, "💰 Unified Billing Engine:", "Auto-aggregates room charges, laboratory tests, medicine dispensation, and doctor fees into a clear, unified invoice.", 13.5, 14)
    add_bullet(tf_hms, "⚡ Modern Tech Stack Architecture:", "React 19 + Node.js cluster architecture, optimized with MySQL 8.0 connection pooling handling hundreds of simultaneous entries.", 13.5, 5)

    # Right Image
    img_path_hms = "c:\\Users\\Hp\\Documents\\HMS\\server_hms\\assets\\hms_dashboard.png"
    if os.path.exists(img_path_hms):
        add_card(s6, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), c_white, c_light_border)
        s6.shapes.add_picture(img_path_hms, Inches(7.0), Inches(1.75), Inches(5.4), Inches(4.7))
    else:
        add_card(s6, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), c_cobalt)

    # =============================================================
    # SLIDE 7: Enterprise Payroll Management System
    # =============================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_background(s7, c_light_bg)
    add_slide_header(s7, "Enterprise Payroll Management System", "Case Study 04")
    
    # Left Specifications
    add_card(s7, Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.0), c_white, c_light_border)
    tf_pay = s7.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.6), Inches(4.4)).text_frame
    tf_pay.word_wrap = True
    
    p_pay_h = tf_pay.paragraphs[0]
    p_pay_h.text = "AUTOMATED SALARY & STAFF RESOURCE MANAGER"
    p_pay_h.font.name = "Segoe UI"
    p_pay_h.font.bold = True
    p_pay_h.font.size = Pt(13)
    p_pay_h.font.color.rgb = c_cobalt
    p_pay_h.space_after = Pt(15)
    
    add_bullet(tf_pay, "🔐 Solid JWT Authentication:", "Secure authentication using JSON Web Tokens (JWT) ensuring secure access for employees and system administrators.", 13.5, 14)
    add_bullet(tf_pay, "👥 Complete Employee Directory:", "Manages employee profiles, work contracts, job roles, attendance logs, and digital tax declaration files.", 13.5, 14)
    add_bullet(tf_pay, "💰 Precision Payroll Calculator:", "Fully automated salary processing factoring in basic pay, custom allowances, tax deductions, unpaid leaves, and bonuses.", 13.5, 14)
    add_bullet(tf_pay, "🗄️ Optimized MongoDB Schemas:", "Employs high-performance NoSQL models for highly flexible logs, historical payroll logs, and dynamic structure alterations.", 13.5, 14)
    add_bullet(tf_pay, "📊 HR Analytics Command Center:", "Sleek charts displaying monthly payouts, salary distributions, active headcount trends, and department expenses.", 13.5, 5)

    # Right Visual representation
    add_card(s7, Inches(7.3), Inches(1.6), Inches(5.2), Inches(5.0), c_dark_navy)
    tf_pay_viz = s7.shapes.add_textbox(Inches(7.6), Inches(2.0), Inches(4.6), Inches(4.2)).text_frame
    tf_pay_viz.word_wrap = True
    p_pay_vt = tf_pay_viz.paragraphs[0]
    p_pay_vt.text = "PAYROLL METRICS CONSOLE"
    p_pay_vt.font.name = "Segoe UI"
    p_pay_vt.font.size = Pt(12)
    p_pay_vt.font.bold = True
    p_pay_vt.font.color.rgb = c_teal
    p_pay_vt.space_after = Pt(25)
    
    p_pay_vb = tf_pay_viz.add_paragraph()
    p_pay_vb.text = "System Integrations:\n\n✔️ Biometric Time Tracker\n✔️ Direct Deposit Bank Generator\n✔️ Custom Payslip PDF Creator\n✔️ Dynamic Deductions Desk\n✔️ Active Staff Dashboard\n\nDatabase: MongoDB Core\nAuth: Secure JWT Session Lock"
    p_pay_vb.font.name = "Segoe UI"
    p_pay_vb.font.size = Pt(16)
    p_pay_vb.font.color.rgb = c_white

    # =============================================================
    # SLIDE 8: RFID Smart Attendance System
    # =============================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_background(s8, c_light_bg)
    add_slide_header(s8, "RFID Smart Attendance System (IoT)", "Case Study 05")
    
    # Left Specs
    add_card(s8, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0), c_white, c_light_border)
    tf_rfid = s8.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.2), Inches(4.4)).text_frame
    tf_rfid.word_wrap = True
    
    p_rfid_h = tf_rfid.paragraphs[0]
    p_rfid_h.text = "HARDWARE-TO-CLOUD REAL-TIME AUTOMATION"
    p_rfid_h.font.name = "Segoe UI"
    p_rfid_h.font.bold = True
    p_rfid_h.font.size = Pt(13)
    p_rfid_h.font.color.rgb = c_teal
    p_rfid_h.space_after = Pt(15)
    
    add_bullet(tf_rfid, "🔌 Arduino Hardware Link:", "Designed real-time serial hardware controller communicating Arduino microcontroller and RFID chip directly to cloud endpoints.", 13.5, 14)
    add_bullet(tf_rfid, "🌐 PHP & SQL Server Engine:", "Robust PHP backend that instantly processes incoming hardware payload, maps RFID serial code, and records entry.", 13.5, 14)
    add_bullet(tf_rfid, "🏆 Innovation Award Winner:", "Formally recognized for engineering innovation by creating a highly cost-efficient IoT hardware prototype.", 13.5, 14)
    add_bullet(tf_rfid, "⚡ Instant Check-in Analytics:", "Live dashboard updating checked-in profiles as cards, highlighting late check-ins, and auto-logging missing members.", 13.5, 14)
    add_bullet(tf_rfid, "📊 Clean Admin Control Portal:", "Responsive, interactive system allowing administrators to easily link RFID tags, manage profiles, and export monthly Excel spreadsheets.", 13.5, 5)

    # Right Image
    img_path_rfid = "c:\\Users\\Hp\\Documents\\HMS\\server_hms\\assets\\rfid_attendance.png"
    if os.path.exists(img_path_rfid):
        add_card(s8, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), c_white, c_light_border)
        s8.shapes.add_picture(img_path_rfid, Inches(7.0), Inches(1.75), Inches(5.4), Inches(4.7))
    else:
        add_card(s8, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), c_cobalt)

    # =============================================================
    # SLIDE 9: Services Offered
    # =============================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_background(s9, c_light_bg)
    add_slide_header(s9, "Core Software Services", "My Offerings")
    
    services = [
        {"title": "Enterprise ERP Development", "desc": "Custom workflows, payroll processing, multi-department academic/clinical inventory automation engineered to fit your processes perfectly."},
        {"title": "Custom SaaS Platforms", "desc": "Scalable, secure multi-tenant web applications with ultra-modern UI layouts, real-time analytics, and secure billing models."},
        {"title": "Interactive Admin Panels", "desc": "Polished developer dashboards built using React, AG Grid, Tailwind, and Recharts, optimized for handling massive data grids seamlessly."},
        {"title": "Robust API Engineering", "desc": "Highly secure, modular RESTful APIs developed with Node.js or PHP, optimized database operations, and micro-second query processing."},
        {"title": "DevOps & Cloud Deploy", "desc": "Professional Linux deployment setups using Docker, Nginx load-balancing, robust SSL integration, and automatic backup routines."},
        {"title": "IoT & System Automation", "desc": "Smart hardware integration (Arduino/RFID microcontrollers) linked directly to web dashboard APIs for real-time monitoring."}
    ]
    
    s_width = Inches(3.6)
    s_height = Inches(2.2)
    s_gap_x = Inches(0.4)
    s_gap_y = Inches(0.4)
    s_start_x = Inches(0.8)
    s_start_y = Inches(1.8)
    
    for idx, srv in enumerate(services):
        row = idx // 3
        col = idx % 3
        
        cx = s_start_x + col * (s_width + s_gap_x)
        cy = s_start_y + row * (s_height + s_gap_y)
        
        # Service card
        add_card(s9, cx, cy, s_width, s_height, c_white, c_light_border)
        
        # Color bar on top of card
        tbar = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, s_width, Inches(0.12))
        tbar.fill.solid()
        tbar.fill.fore_color.rgb = c_cobalt if idx % 2 == 0 else c_teal
        tbar.line.fill.background()
        
        # Content box
        srv_box = s9.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.2), s_width - Inches(0.3), s_height - Inches(0.35))
        tf_srv = srv_box.text_frame
        tf_srv.word_wrap = True
        
        p_title = tf_srv.paragraphs[0]
        p_title.text = srv["title"]
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark_navy
        p_title.space_after = Pt(8)
        
        p_desc = tf_srv.add_paragraph()
        p_desc.text = srv["desc"]
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = c_text_muted

    # =============================================================
    # SLIDE 10: Why Choose Me
    # =============================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_background(s10, c_light_bg)
    add_slide_header(s10, "Why Partner With Me?", "My Value Proposition")
    
    reasons = [
        {"title": "💼 Battle-Tested Real World Experience", "desc": "Proven track record designing heavy administrative systems handling tens of thousands of real active student registrations and medical records daily."},
        {"title": "⚡ Enterprise-Class Resilient Architecture", "desc": "Fluent in implementing Node.js cluster modes, database connection pooling, parameterized indexing, and optimal JWT/RBAC security protocols."},
        {"title": "🎨 Exceptional Custom UI/UX & Responsive Flow", "desc": "Clean, pixel-perfect modern web layouts engineered using React 19, AG Grid, Framer Motion, and Tailwind CSS. No generic, boring placeholders used."},
        {"title": "🚀 Professional Full-Cycle System Delivery", "desc": "Complete project execution starting from requirements mapping, wireframing, core programming, staging tests, down to robust live server setups."}
    ]
    
    r_width = Inches(5.6)
    r_height = Inches(2.2)
    r_gap_x = Inches(0.5)
    r_gap_y = Inches(0.4)
    r_start_x = Inches(0.8)
    r_start_y = Inches(1.8)
    
    for idx, rsn in enumerate(reasons):
        row = idx // 2
        col = idx % 2
        
        cx = r_start_x + col * (r_width + r_gap_x)
        cy = r_start_y + row * (r_height + r_gap_y)
        
        add_card(s10, cx, cy, r_width, r_height, c_white, c_light_border)
        
        # Elegant side strip accent
        lstrip = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.12), r_height)
        lstrip.fill.solid()
        lstrip.fill.fore_color.rgb = c_teal if idx % 2 == 0 else c_cobalt
        lstrip.line.fill.background()
        
        rsn_box = s10.shapes.add_textbox(cx + Inches(0.3), cy + Inches(0.2), r_width - Inches(0.45), r_height - Inches(0.4))
        tf_rsn = rsn_box.text_frame
        tf_rsn.word_wrap = True
        
        p_title = tf_rsn.paragraphs[0]
        p_title.text = rsn["title"]
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark_navy
        p_title.space_after = Pt(10)
        
        p_desc = tf_rsn.add_paragraph()
        p_desc.text = rsn["desc"]
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = c_text_muted

    # =============================================================
    # SLIDE 11: Development Process / Workflow
    # =============================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_background(s11, c_light_bg)
    add_slide_header(s11, "Professional Development Process", "My Workflow")
    
    workflow = [
        {"num": "01", "step": "Requirements Mapping", "desc": "Mapping business specifications, data schemas, and key integrations."},
        {"num": "02", "step": "Sleek UI/UX Wireframing", "desc": "Prototyping premium responsive interfaces focused on ease of operation."},
        {"num": "03", "step": "Agile Development Cycle", "desc": "Structured clean coding using modern backend clusters & React modules."},
        {"num": "04", "step": "Rigorous End-to-End Test", "desc": "Automated security vetting, database query performance profiling."},
        {"num": "05", "step": "Production Deployment", "desc": "Linux server configuration with Nginx reverse proxy, SSL secure locks."},
        {"num": "06", "step": "Dedicated Aftercare", "desc": "Active system health telemetry monitoring and performance updates."}
    ]
    
    w_width = Inches(1.8)
    w_height = Inches(4.5)
    w_gap = Inches(0.18)
    w_start_x = Inches(0.8)
    w_start_y = Inches(1.8)
    
    for idx, step in enumerate(workflow):
        cx = w_start_x + idx * (w_width + w_gap)
        
        # Step card
        add_card(s11, cx, w_start_y, w_width, w_height, c_white, c_light_border)
        
        # Header numbering card
        num_shape = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, w_start_y, w_width, Inches(0.7))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = c_dark_navy if idx % 2 == 0 else c_cobalt
        num_shape.line.fill.background()
        
        tf_num = num_shape.text_frame
        tf_num.word_wrap = True
        p_num = tf_num.paragraphs[0]
        p_num.text = step["num"]
        p_num.font.name = "Segoe UI"
        p_num.font.size = Pt(18)
        p_num.font.bold = True
        p_num.font.color.rgb = c_teal
        p_num.alignment = PP_ALIGN.CENTER
        
        # Text details
        st_box = s11.shapes.add_textbox(cx + Inches(0.1), w_start_y + Inches(0.8), w_width - Inches(0.2), w_height - Inches(0.9))
        tf_st = st_box.text_frame
        tf_st.word_wrap = True
        
        p_title = tf_st.paragraphs[0]
        p_title.text = step["step"]
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark_navy
        p_title.space_after = Pt(12)
        
        p_desc = tf_st.add_paragraph()
        p_desc.text = step["desc"]
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = c_text_muted

    # =============================================================
    # SLIDE 12: Contact Information (Dark Theme)
    # =============================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_background(s12, c_dark_navy)
    
    # Glowing bottom strip
    g12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.25), Inches(13.333), Inches(0.25))
    g12.fill.solid()
    g12.fill.fore_color.rgb = c_teal
    g12.line.fill.background()
    
    # Call to action
    cta_box = s12.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.4))
    cta_box.text_frame.word_wrap = True
    p_cta = cta_box.text_frame.paragraphs[0]
    p_cta.text = "LET'S BUILD SOMETHING EXCEPTIONAL"
    p_cta.font.name = "Segoe UI"
    p_cta.font.size = Pt(12)
    p_cta.font.bold = True
    p_cta.font.color.rgb = c_teal
    
    # Heading
    end_box = s12.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(1.2))
    tf_end = end_box.text_frame
    tf_end.word_wrap = True
    p_end = tf_end.paragraphs[0]
    p_end.text = "Partner With an Enterprise\nSoftware Expert Today."
    p_end.font.name = "Segoe UI"
    p_end.font.size = Pt(36)
    p_end.font.bold = True
    p_end.font.color.rgb = c_white
    
    # Grid of contact cards
    channels = [
        {"title": "📧 EMAIL CHANNEL", "val": "prawin1209@gmail.com", "sub": "Direct Communication"},
        {"title": "🌐 DIGITAL PORTFOLIO", "val": "prawin.vercel.app", "sub": "Live Demos & Projects"},
        {"title": "🐙 GITHUB SOURCE", "val": "github.com/prawinkumar2k", "sub": "Clean Repositories & Diffs"},
        {"title": "👥 PROFESSIONAL LINKEDIN", "val": "linkedin.com/in/prawinkumar-n", "sub": "Industry Networking"}
    ]
    
    ch_width = Inches(5.4)
    ch_height = Inches(1.4)
    ch_gap_x = Inches(0.5)
    ch_gap_y = Inches(0.4)
    ch_start_x = Inches(1.0)
    ch_start_y = Inches(3.6)
    
    for idx, ch in enumerate(channels):
        row = idx // 2
        col = idx % 2
        
        cx = ch_start_x + col * (ch_width + ch_gap_x)
        cy = ch_start_y + row * (ch_height + ch_gap_y)
        
        # Dark cobalt card
        add_card(s12, cx, cy, ch_width, ch_height, c_cobalt)
        
        ch_box = s12.shapes.add_textbox(cx + Inches(0.25), cy + Inches(0.25), ch_width - Inches(0.5), ch_height - Inches(0.5))
        tf_ch = ch_box.text_frame
        tf_ch.word_wrap = True
        
        p_ch_t = tf_ch.paragraphs[0]
        p_ch_t.text = ch["title"]
        p_ch_t.font.name = "Segoe UI"
        p_ch_t.font.size = Pt(11)
        p_ch_t.font.bold = True
        p_ch_t.font.color.rgb = c_teal
        p_ch_t.space_after = Pt(4)
        
        p_ch_v = tf_ch.add_paragraph()
        p_ch_v.text = ch["val"]
        p_ch_v.font.name = "Segoe UI"
        p_ch_v.font.size = Pt(16)
        p_ch_v.font.bold = True
        p_ch_v.font.color.rgb = c_white
        p_ch_v.space_after = Pt(2)
        
        p_ch_s = tf_ch.add_paragraph()
        p_ch_s.text = ch["sub"]
        p_ch_s.font.name = "Segoe UI"
        p_ch_s.font.size = Pt(10)
        p_ch_s.font.color.rgb = RGBColor(165, 180, 252)

    # Save presentation
    output_filename = "c:\\Users\\Hp\\Documents\\HMS\\server_hms\\assets\\PrawinKumar_Developer_Portfolio.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully created and saved at: {output_filename}")

if __name__ == "__main__":
    create_presentation()
